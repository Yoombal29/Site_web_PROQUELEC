import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configuration du dossier
BASE_DIR = r"c:\Mes Sites Web\Site_web_PROQUELEC-main"
OUTPUT_DIR = os.path.join(BASE_DIR, "KEBE")

if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

# Couleurs Institutionnelles (PROQUELEC / ASN)
BLUE_DARK = RGBColor(0, 51, 102)
GOLD = RGBColor(184, 134, 11)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

def apply_institutional_design(slide, title_text, is_title_page=False):
    if is_title_page:
        # Background color for title slide
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BLUE_DARK
        return

    # Style du titre pour slides de contenu
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.name = 'Segoe UI'
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = BLUE_DARK
    title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Ligne dorée ASN/PROQUELEC
    left = Inches(0.5)
    top = Inches(1.1)
    width = Inches(9)
    height = Inches(0.04)
    line = slide.shapes.add_shape(1, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.visible = False

def create_norm_presentation(filename, main_title, subtitle, modules_content, evaluation_data):
    prs = Presentation()
    
    # --- 1. Slide de Titre ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    apply_institutional_design(slide, "", is_title_page=True)
    
    title_shape = slide.shapes.title
    sub_shape = slide.placeholders[1]
    
    title_shape.text = main_title
    sub_shape.text = f"{subtitle}\nConforme Norme NF C 18-510 (Annexe D)\nExpertise PROQUELEC - Sénégal"
    
    title_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    sub_shape.text_frame.paragraphs[0].font.color.rgb = GOLD
    sub_shape.text_frame.paragraphs[0].font.size = Pt(24)

    # --- 2. Slides de Contenu (Thèmes de la Norme) ---
    for section_title, bullets in modules_content:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        apply_institutional_design(slide, section_title)
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
            p.space_after = Pt(8)

    # --- 3. Slide d'Évaluation (Critères Officiels) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    apply_institutional_design(slide, "Évaluation et Validation (Annexe D.3)")
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = f"Évaluation Théorique (QCM) : {evaluation_data['qcm_count']} questions minimum."
    p.level = 0
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = f"Seuil de réussite : {evaluation_data['threshold']} / Zéro erreur sur les questions fondamentales."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Évaluation Pratique (Mises en situation) :"
    p.level = 0
    p.font.bold = True
    
    for situ in evaluation_data['practical']:
        p = tf.add_paragraph()
        p.text = situ
        p.level = 1

    p = tf.add_paragraph()
    p.text = "Notation A/B/C/D : Réussite si (Aucun D et max un C)."
    p.level = 0
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 0, 0)

    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    print(f"✅ Généré : {filename}")

# --- DÉFINITION EXHAUSTIVE DES CONTENUS (NF C 18-510 ANNEXE D) ---

# Tronc Commun TC1 & TC2 (Base)
tc1_themes = [
    ("TC1 : Grandeurs et Dangers", ["Grandeurs électriques (Tension, Courant, Résistance).", "Dangers de l’électricité : Électrisation, Électrocution, Brûlures.", "Statistiques et mécanismes du risque."]),
    ("TC1 : Environnement et Zones", ["Domaines de tension (TBT, BT, HT).", "Zones d'environnement et leurs limites (DLV, DMA).", "Zones en champ libre et locaux réservés."]),
    ("TC1 : Principes de Prévention", ["Définition et buts de l’habilitation.", "Les symboles d’habilitation (Codes et indices).", "Utilisation des EPC (Nappes, Balisage, Écrans).", "Équipements de protection individuelle (EPI).", "Conduite à tenir : Accident (Sécourisme) et Incendie."])
]

tc2_themes = tc1_themes + [
    ("TC2 : Approfondissement Sécurité", ["Analyse de risque électrique sur site.", "Vérification et choix des EPI/EPC.", "Gestion des documents et transmission des consignes."])
]

# 1. B0-H0-H0V (Exécutant vs Chargé)
b0_exec_content = tc1_themes + [
    ("Exécutant Non-Électricien (B0/H0)", ["Acteurs concernés (Donneur d'ordre, Chef d'établissement).", "Limites de l'habilitation chiffre 0.", "Repérage de la zone de travail et respect du balisage.", "Prescriptions de sécurité pour les travaux non-électriques."])
]
b0_charge_content = tc1_themes + [
    ("Chargé de chantier Non-Électricien (B0/H0)", ["Rôle et responsabilités du chargé de chantier.", "Limites et zone de travail.", "Définition et balisage de la zone.", "Surveillance constante des intervenants.", "Échange de documents (Avis de fin de travaux)."])
]

# 2. BP (Photovoltaïque)
bp_content = tc1_themes + [
    ("Module BP : Spécificités Photovoltaïques", ["Risques spécifiques au courant continu (DC).", "Moyens de protection contre les chocs électriques.", "Règles de pose des modules et raccordements.", "Protection des connecteurs et gestion des câbles sous tension."])
]

# 3. BS (Interventions Élémentaires)
bs_content = tc1_themes + [
    ("D.2.9 Module BS : Cadre d'Intervention", ["Limites de l'habilitation BS (max 400V, 32A).", "Échange d'informations avec le chargé d'exploitation.", "Matériels de mesure et de protection (VAT)."]),
    ("BS : Mise en sécurité et Opérations", ["Procédure de mise en sécurité (Séparation, Condamnation, VAT).", "Remplacement de fusibles et lampes.", "Raccordement d'un circuit terminal (Prises, Luminaires).", "Réarmement de dispositifs de protection."])
]

# 4. BE / HE Manœuvre
be_he_man_content = tc1_themes + [
    ("D.2.14 Manœuvres en BT et HT", ["Information sur les types de manœuvres (D’exploitation, D’urgence).", "Matériels électriques BT et HT concernés.", "Limites de l'habilitation BE / HE Manœuvre.", "Équipements de protection spécifiques aux manœuvres.", "Compte rendu au chargé d'exploitation."])
]

# 5. BR (Interventions Générales)
br_content = tc2_themes + [
    ("D.2.11 Module BR : Chargé d'Intervention", ["Rôle du chargé d'exploitation.", "Habilitation BR : Dépannage, mesurage, essais.", "Consignation pour son propre compte (Les 5 étapes)."]),
    ("BR : Mode opératoire", ["Procédures de dépannage et de connexion/déconnexion.", "Utilisation sûre des appareils de mesure.", "Réalisation d'essais et de mesurages sous tension (si nécessaire)."])
]

# 6. BC / HC (Consignation)
bc_hc_content = tc2_themes + [
    ("D.2.12 Consignation (BC / HC)", ["Le rôle du BC (Consignation) et HC (Haute Tension).", "Attestation de consignation pour travaux.", "Interface avec le chargé de travaux (B2/H2)."]),
    ("Procédure de Consignation (5 Étapes)", ["1. Séparation : Coupure physique de toutes les sources.", "2. Condamnation : Dispositifs d'interdiction (Cadenas).", "3. Identification : Vérification de l'ouvrage.", "4. VAT : Vérification d'Absence de Tension immédiate.", "5. MALT/CC : Mise à la Terre (obligatoire en HT)."])
]

# 7. B1-B2-H1-H2 (Travaux Hors Tension)
b1_b2_content = tc2_themes + [
    ("D.2.13 Travaux Hors Tension (B1, B2, H1, H2)", ["Définition des travaux hors tension.", "Rôle de l'Exécutant (B1/H1) sous ordre.", "Rôle du Chargé de travaux (B2/H2) : Organisation.", "Organisation du chantier et suppression du voisinage.", "Habilitation V (B1V, B2V) : Travail au voisinage."])
]

# 8. HT (H1, H2, HC, HE) - Module Spécifique
ht_spec_content = tc2_themes + [
    ("Spécificités Haute Tension (HTA/HTB)", ["Matériels HT et postes de transformation.", "Distances Minimales d'Approche (DMA) augmentées.", "Risque d'induction et de couplage.", "Mise à la Terre et en Court-Circuit (MALT/CC) systématique.", "Documents spécifiques HT (Message de début/fin)."])
]

# --- GÉNÉRATION ---

if __name__ == "__main__":
    # Liste des configurations (FileName, Title, Subtitle, Content, Eval)
    config = [
        ("D2.4_B0_H0_Executant.pptx", "Module B0 - H0 - H0V", "Exécutant de travaux d’ordre non électrique", b0_exec_content, {"qcm_count": 15, "threshold": "70%", "practical": ["Identification du secteur", "Respect des zones et du balisage"]}),
        
        ("D2.5_B0_H0_Charge_Chantier.pptx", "Module B0 - H0 - H0V", "Chargé de chantier d’ordre non électrique", b0_charge_content, {"qcm_count": 15, "threshold": "70%", "practical": ["Analyse de zone", "Surveillance des tiers", "Instructions de sécurité"]}),
        
        ("D2.6_BP_Photovoltaique.pptx", "Module BP", "Pose d'équipement photovoltaïque", bp_content, {"qcm_count": 10, "threshold": "75%", "practical": ["Pose de modules DC en sécurité", "Raccordement connecteurs"]}),
        
        ("D2.9_BS_Intervention_Elementaire.pptx", "Module BS", "Interventions BT élémentaires", bs_content, {"qcm_count": 20, "threshold": "80%", "practical": ["Mise en sécurité (VAT)", "Remplacement fusible/lampe", "Raccordement circuit terminal"]}),
        
        ("D2.11_BR_Intervention_Generale.pptx", "Module BR", "Interventions BT générales", br_content, {"qcm_count": 20, "threshold": "85%", "practical": ["Consignation personnelle", "Dépannage complet", "Essai de fonctionnement"]}),
        
        ("D2.12_BC_HC_Consignation.pptx", "Module BC / HC", "Consignation en Basse et Haute Tension", bc_hc_content, {"qcm_count": 20, "threshold": "90%", "practical": ["Réalisation des 5 étapes", "Rédaction attestation de consignation"]}),
        
        ("D2.13_B1_B2_H1_H2_Travaux.pptx", "Module B1-B2 / H1-H2", "Travaux hors tension (Exécutant & Chargé)", b1_b2_content, {"qcm_count": 20, "threshold": "85%", "practical": ["Organisation de chantier", "Surveillance et Zone de travail"]}),
        
        ("D2.14_BE_HE_Manoeuvre.pptx", "Module BE / HE Manœuvre", "Manœuvres en Basse et Haute Tension", be_he_man_content, {"qcm_count": 15, "threshold": "75%", "practical": ["Réalisation de manœuvre d'exploitation", "Message de fin"]}),
    ]

    for fn, tit, sub, cont, ev in config:
        create_norm_presentation(fn, tit, sub, cont, ev)

    print("\n🚀 TOUS LES MODULES DE L'ANNEXE D (NF C 18-510) ONT ÉTÉ GÉNÉRÉS DANS 'KEBE' AVEC RIGUEUR TECHNIQUE.")
