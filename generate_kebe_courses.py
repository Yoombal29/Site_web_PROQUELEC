import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configuration du dossier
BASE_DIR = r"c:\Mes Sites Web\Site_web_PROQUELEC-main"
OUTPUT_DIR = os.path.join(BASE_DIR, "KEBE")

if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

# Couleurs PROQUELEC
BLUE_DARK = RGBColor(0, 32, 96)
GOLD = RGBColor(197, 160, 89)
WHITE = RGBColor(255, 255, 255)

def apply_clean_design(slide, title_text):
    # Fond blanc (par défaut)
    # Titre avec barre bleue
    title_shape = slide.shapes.title
    title_shape.text = title_text
    
    # Style du titre
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.name = 'Arial'
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = BLUE_DARK
    title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Ligne de séparation Gold
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(0.05)
    line = slide.shapes.add_shape(1, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.visible = False

def create_presentation(filename, title, slides_content):
    prs = Presentation()
    
    # Slide de titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title_shape.text = title
    subtitle.text = "Formation à l'Habilitation Électrique - NF C 18-510\nExpert Formateur PROQUELEC"
    
    title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # Slides de contenu
    for s_title, s_bullets in slides_content:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        apply_clean_design(slide, s_title)
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for bullet in s_bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)
            p.font.name = 'Calibri'
            p.space_after = Pt(10)

    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    print(f"✅ Présentation générée : {filename}")

# --- MODULES DE FORMATION ---

# 1. B0 - H0 - H0V
module_b0 = [
    ("Le Risque Électrique", ["Analyses statistiques et causes d'accidents.", "Les effets du courant sur le corps (Électrisation vs Électrocution).", "Notions de tension (BT / HT)."]),
    ("L'Environnement Électrique", ["Les zones d'application (Voisinage, Champ libre).", "Distances Limites de Voisinage (DLV).", "Distances Minimales d'Approche (DMA)."]),
    ("L'Habilitation Électrique", ["Définition : Reconnaissance de la capacité à travailler en sécurité.", "Symboles d'habilitation : Qui fait quoi ?", "Le Titre d'Habilitation : Validité et limites."]),
    ("Prévention et Protection", ["Les EPI (Gants isolants, Écran facial).", "Les EPC (Nappes, Balisage, Obstacles).", "Consignes de sécurité pour non-électriciens."])
]

# 2. BS - Intervention Élémentaire
module_bs = module_b0 + [
    ("Le Module BS", ["Remplacement de fusibles, lampes et accessoires.", "Raccordement d'un circuit terminal (max 400V, 32A).", "Réarmement de protection."]),
    ("Procédure de Mise en Sécurité", ["Séparation du circuit.", "Condamnation (Cadenas et MAC).", "VAT (Vérification d'Absence de Tension)."]),
    ("Documents et Avis", ["Le bon d'intervention.", "Avis de fin d'intervention."])
]

# 3. BC - Chargé de Consignation
module_bc = module_b0 + [
    ("Le Rôle du BC", ["Responsabilité de la mise en sécurité.", "Interface entre le donneur d'ordre et le chargé de travaux."]),
    ("Les 5 étapes de la Consignation", ["1. Séparation : Coupure physique.", "2. Condamnation : Interdiction de manœuvre.", "3. Identification : Vérification de l'ouvrage.", "4. VAT : Vérification d'Absence de Tension.", "5. MALT/CC : Mise à la terre et en court-circuit."]),
    ("Documents de Consignation", ["Attestation de Consignation.", "Fiche de manœuvre."])
]

# 4. BR - Chargé d'Intervention Générale
module_br = module_b0 + [
    ("Spécificités du BR", ["Dépannage, maintenance et mesurage.", "Périmètre d'intervention limité (BT)."]),
    ("Mesurages et Essais", ["Utilisation du multimètre et de la pince ampèremétrique.", "Risques liés aux pointes de touche.", "Essais de fonctionnement sous tension."]),
    ("Règle d'Or", ["Travail hors tension prioritaire.", "Mesures compensatoires si travail au voisinage."])
]

# 5. B1 / B2 - Travaux Électriques
module_b1b2 = module_b0 + [
    ("Exécutants et Chargés de Travaux", ["B1 : Exécutant (Obéit au B2).", "B2 : Organise et surveille le chantier."]),
    ("Le Voisinage", ["Habilitation V (B1V, B2V).", "Port des EPI obligatoire dès l'entrée en DLV."]),
    ("Organisation du Chantier", ["Préparation de la zone.", "Surveillance constante des intervenants."])
]

if __name__ == "__main__":
    create_presentation("Module_B0_H0_Travaux_Non_Electriques.pptx", "Formation Habilitation B0 / H0 / H0V", module_b0)
    create_presentation("Module_BS_Intervention_Elementaire.pptx", "Formation Habilitation BS", module_bs)
    create_presentation("Module_BC_Charge_de_Consignation.pptx", "Formation Habilitation BC", module_bc)
    create_presentation("Module_BR_Charge_d_Intervention_Generale.pptx", "Formation Habilitation BR", module_br)
    create_presentation("Module_B1_B2_Travaux_Electriques.pptx", "Formation Habilitation B1 / B2", module_b1b2)
